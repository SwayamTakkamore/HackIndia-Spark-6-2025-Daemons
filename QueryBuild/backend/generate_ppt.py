from pptx import Presentation
from pptx.util import Inches, Pt
from .llm_wrapper import query_llm

def format_for_slide(text, max_bullets=5):
    """Format and limit bullet points to prevent overflow"""
    # Split by bullet points (both • and - are common bullet indicators)
    if '•' in text:
        bullets = [b.strip() for b in text.split('•') if b.strip()]
    elif '-' in text:
        bullets = [b.strip() for b in text.split('-') if b.strip()]
    else:
        # If no bullet indicators, just split by newlines
        bullets = [b.strip() for b in text.split('\n') if b.strip()]
    
    # Format each bullet point and limit length to ~100 chars
    formatted_bullets = []
    for bullet in bullets[:max_bullets]:
        if len(bullet) > 100:
            bullet = bullet[:97] + "..."
        formatted_bullets.append(f"• {bullet}")
    
    return "\n".join(formatted_bullets)

def add_content_slide(prs, title, content):
    """Add a slide with proper text formatting"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    # Format and add content
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()  # Clear existing placeholder text
    
    # Set proper font size
    formatted_content = format_for_slide(content)
    p = text_frame.paragraphs[0]
    p.text = formatted_content
    
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)  # Adjust font size to fit slide

def generate_ppt(problem_statement, save_path):
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project Presentation"
    subtitle.text = problem_statement[:150] + ("..." if len(problem_statement) > 150 else "")

    # Introduction slide
    intro_prompt = f"Given the problem statement: '{problem_statement}', provide a brief introduction and background of the problem in 3-4 bullet points."
    intro_points = query_llm(intro_prompt)
    add_content_slide(prs, "Introduction", intro_points)
    
    # Methodology slide
    method_prompt = f"Given the problem statement: '{problem_statement}', outline the methodology and approach to solve this in 4-5 bullet points."
    method_points = query_llm(method_prompt)
    add_content_slide(prs, "Methodology", method_points)
    
    # Results slide
    results_prompt = f"Given the problem statement: '{problem_statement}', what results and insights might we expect? Provide 3-4 bullet points."
    results_points = query_llm(results_prompt)
    add_content_slide(prs, "Expected Results", results_points)
    
    # Conclusion slide
    add_content_slide(prs, "Conclusion", "• Summary of approach\n• Key findings\n• Future work")
    
    prs.save(save_path)
    return save_path
